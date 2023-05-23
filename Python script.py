from docx import Document
from pptx import Presentation

def copy_doc_to_ppt(doc_path, ppt_path):
    # Load the DOC file
    doc = Document(doc_path)

    # Create a new PPT presentation
    presentation = Presentation()

    # Iterate through the paragraphs in the DOC file
    for paragraph in doc.paragraphs:
        # Create a new slide
        slide_layout = presentation.slide_layouts[1]  # Choose the slide layout (e.g., Title and Content)
        slide = presentation.slides.add_slide(slide_layout)

        # Add the paragraph content to the slide
        content = slide.shapes.add_textbox().text_frame
        content.text = paragraph.text

    # Save the PPT presentation
    presentation.save(ppt_path)
    print("Content copied from DOC to PPT successfully!")

# Provide the file paths for the input DOC and output PPT files
doc_file_path = "path/to/input.docx"
ppt_file_path = "path/to/output.pptx"

# Call the function to copy the content from DOC to PPT
copy_doc_to_ppt(doc_file_path, ppt_file_path)
